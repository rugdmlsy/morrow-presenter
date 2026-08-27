#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=1.0.2", "Pillow>=10"]
# ///
from __future__ import annotations
import hashlib,json,sys,uuid
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.dml import MSO_FILL_TYPE

SLIDE_ASPECT=16/9

def uid(): return str(uuid.uuid4())
def pct(v,total): return float(v)/float(total)*100 if total else 0.0
def color_of(fmt, default):
    try:
        if fmt.type==MSO_FILL_TYPE.SOLID and fmt.fore_color.rgb is not None: return '#'+str(fmt.fore_color.rgb)
    except Exception: pass
    return default

def line_color(shape,default='#4b4d50'):
    try:
        if shape.line.fill.type==MSO_FILL_TYPE.SOLID and shape.line.color.rgb is not None:return '#'+str(shape.line.color.rgb)
    except Exception: pass
    return default

def line_width(shape,default=1.5):
    try:return max(.0,float(shape.line.width.pt)) if shape.line.width is not None else default
    except Exception:return default

def first_text_style(shape):
    out={'fontFamily':'Inter','fontSize':24.0,'fontWeight':400,'italic':False,'underline':False,'color':'#202124','align':'left','verticalAlign':'middle'}
    try:
        tf=shape.text_frame
        if tf.paragraphs:
            p=tf.paragraphs[0]
            amap={1:'left',2:'center',3:'right'}
            if p.alignment is not None:out['align']=amap.get(int(p.alignment),'left')
            if p.runs:
                f=p.runs[0].font
                if f.name:out['fontFamily']=f.name
                if f.size:out['fontSize']=float(f.size.pt)
                if f.bold is not None:out['fontWeight']=700 if f.bold else 400
                if f.italic is not None:out['italic']=bool(f.italic)
                if f.underline is not None:out['underline']=bool(f.underline)
                try:
                    if f.color.rgb is not None:out['color']='#'+str(f.color.rgb)
                except Exception:pass
        va=getattr(tf,'vertical_anchor',None)
        if va is not None:out['verticalAlign']={1:'top',3:'middle',4:'bottom'}.get(int(va),'middle')
    except Exception:pass
    return out

def common(shape,sw,sh,group_id=None):
    return {'id':uid(),'x':pct(shape.left,sw),'y':pct(shape.top,sh),'width':max(.5,pct(shape.width,sw)),'height':max(.5,pct(shape.height,sh)),'rotation':float(getattr(shape,'rotation',0) or 0)%360,'opacity':1.0,'locked':False,'groupId':group_id}

def asset_from_picture(shape,outdeck):
    image=shape.image; blob=image.blob; ext=(getattr(image,'ext',None) or 'png').lower(); digest=hashlib.sha256(blob).hexdigest()[:24]; rel=f'.morrow-assets/{digest}.{ext}'; target=outdeck.parent/rel;target.parent.mkdir(parents=True,exist_ok=True);target.write_bytes(blob) if not target.exists() else None
    try:
        with Image.open(target) as im:iw,ih=im.size
    except Exception:iw,ih=16,9
    return rel,float(iw),float(ih)

def picture_element(shape,sw,sh,outdeck,group_id=None):
    rel,iw,ih=asset_from_picture(shape,outdeck);visx=pct(shape.left,sw);visy=pct(shape.top,sh);visw=pct(shape.width,sw);vish=pct(shape.height,sh)
    def crop(name):
        try:return max(0.0,min(.95,float(getattr(shape,name) or 0)))
        except:return 0.0
    l,t,r,b=(crop('crop_left'),crop('crop_top'),crop('crop_right'),crop('crop_bottom')); fw=visw/max(.05,1-l-r); fh=vish/max(.05,1-t-b); expected=fw*SLIDE_ASPECT/(iw/ih if iw and ih else SLIDE_ASPECT)
    if abs(expected-fh)>max(1.5,.08*max(expected,fh)):
        fw=fh*(iw/ih if iw and ih else SLIDE_ASPECT)/SLIDE_ASPECT; expected=fh
    return {'id':uid(),'type':'image','path':rel,'alt':getattr(shape,'name',''),'x':visx-fw*l,'y':visy-expected*t,'width':fw,'height':expected,'intrinsicWidth':iw,'intrinsicHeight':ih,'crop':{'left':l*100,'top':t*100,'right':r*100,'bottom':b*100},'rotation':float(getattr(shape,'rotation',0) or 0)%360,'opacity':1.0,'locked':False,'groupId':group_id}

def placeholder_role(shape):
    try:
        typ=shape.placeholder_format.type
        if typ in (PP_PLACEHOLDER.TITLE,PP_PLACEHOLDER.CENTER_TITLE):return 'title'
        if typ in (PP_PLACEHOLDER.BODY,PP_PLACEHOLDER.SUBTITLE):return 'body'
    except Exception:pass
    return None

def text_element(shape,sw,sh,group_id=None,role=None):
    if role is None:
        name=getattr(shape,'name','') or ''
        if name.startswith('MorrowPresenter:') and name.split(':',1)[1] in ('title','body'): role=name.split(':',1)[1]
    c=common(shape,sw,sh,group_id);style=first_text_style(shape);fill='transparent';stroke='transparent';swid=0.0
    try:fill=color_of(shape.fill,'transparent')
    except:pass
    try:stroke=line_color(shape,'transparent');swid=line_width(shape,0.0)
    except:pass
    return {**c,'type':'text','role':role,'text':shape.text or '','fill':fill,'stroke':stroke,'strokeWidth':swid,'padding':1.2,**style}

def shape_kind(shape):
    try:
        a=shape.auto_shape_type
        if a==MSO_SHAPE.OVAL:return'ellipse'
        if a==MSO_SHAPE.ROUNDED_RECTANGLE:return'rounded-rect'
        if a==MSO_SHAPE.RECTANGLE:return'rect'
    except:pass
    return'rect'

def auto_shape_element(shape,sw,sh,group_id=None):
    c=common(shape,sw,sh,group_id);style=first_text_style(shape);return {**c,'type':'shape','shape':shape_kind(shape),'text':getattr(shape,'text','') or '','fill':color_of(shape.fill,'#e8e8e4'),'stroke':line_color(shape),'strokeWidth':line_width(shape),**style}

def line_element(shape,sw,sh,group_id=None):
    c=common(shape,sw,sh,group_id);return {**c,'type':'shape','shape':'line','text':'','fill':'transparent','stroke':line_color(shape),'strokeWidth':line_width(shape),'fontFamily':'Inter','fontSize':20,'fontWeight':400,'italic':False,'underline':False,'color':'#202124','align':'center','verticalAlign':'middle'}

def table_element(shape,sw,sh,group_id=None):
    c=common(shape,sw,sh,group_id);t=shape.table;rows=len(t.rows);cols=len(t.columns);cells=[[t.cell(r,cc).text or '' for cc in range(cols)] for r in range(rows)];style={'fontFamily':'Inter','fontSize':18,'fontWeight':400,'italic':False,'underline':False,'color':'#202124','align':'left','verticalAlign':'middle'}
    try:
        cell=t.cell(0,0);tf=cell.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            f=tf.paragraphs[0].runs[0].font
            if f.name:style['fontFamily']=f.name
            if f.size:style['fontSize']=float(f.size.pt)
    except:pass
    return {**c,'type':'table','rows':rows,'cols':cols,'cells':cells,'fill':'#ffffff','headerFill':'#e8e8e4','stroke':'#777777','strokeWidth':1.0,**style}

def fallback_box(shape,sw,sh,label,group_id=None):
    c=common(shape,sw,sh,group_id);return {**c,'type':'shape','shape':'rect','text':label,'fill':'#f4f4f2','stroke':'#b8b8b3','strokeWidth':1.0,'fontFamily':'Inter','fontSize':16,'fontWeight':400,'italic':False,'underline':False,'color':'#666666','align':'center','verticalAlign':'middle'}

def walk(shapes,sw,sh,outdeck,group_id=None):
    out=[]
    for shape in shapes:
        try:stype=shape.shape_type
        except:continue
        if stype==MSO_SHAPE_TYPE.GROUP:
            gid=uid();out.extend(walk(shape.shapes,sw,sh,outdeck,gid));continue
        try:
            if stype==MSO_SHAPE_TYPE.PICTURE:out.append(picture_element(shape,sw,sh,outdeck,group_id));continue
            if getattr(shape,'has_table',False):out.append(table_element(shape,sw,sh,group_id));continue
            if stype==MSO_SHAPE_TYPE.LINE:out.append(line_element(shape,sw,sh,group_id));continue
            if getattr(shape,'has_chart',False):out.append(fallback_box(shape,sw,sh,'[Chart]',group_id));continue
            if getattr(shape,'has_text_frame',False):
                if stype==MSO_SHAPE_TYPE.PLACEHOLDER:out.append(text_element(shape,sw,sh,group_id,placeholder_role(shape)))
                elif stype==MSO_SHAPE_TYPE.AUTO_SHAPE:out.append(auto_shape_element(shape,sw,sh,group_id))
                else:out.append(text_element(shape,sw,sh,group_id))
                continue
            out.append(fallback_box(shape,sw,sh,f'[{getattr(shape,"name","Object")}]',group_id))
        except Exception as exc:
            out.append(fallback_box(shape,sw,sh,f'[Unsupported: {getattr(shape,"name","Object")}]',group_id))
    return out

def slide_background(slide):
    try:return color_of(slide.background.fill,'#ffffff')
    except:return'#ffffff'

def notes(slide):
    try:return slide.notes_slide.notes_text_frame.text or ''
    except:return''

def main():
    if len(sys.argv)!=3:raise SystemExit('usage: import-pptx.py input.pptx output.morrowdeck')
    source=Path(sys.argv[1]).resolve();outdeck=Path(sys.argv[2]).resolve();prs=Presentation(source);sw,sh=prs.slide_width,prs.slide_height;slides=[]
    for slide in prs.slides:
        slides.append({'id':uid(),'layout':'blank','background':slide_background(slide),'notes':notes(slide),'transition':{'type':'none','duration':.35},'elements':walk(slide.shapes,sw,sh,outdeck)})
    if not slides:slides=[{'id':uid(),'layout':'blank','background':'#ffffff','notes':'','transition':{'type':'none','duration':.35},'elements':[]}]
    deck={'version':1,'title':source.stem,'selectedId':slides[0]['id'],'theme':{'name':'default','fontFamily':'Inter','titleFontFamily':'Inter','background':'#ffffff','text':'#202124','accent':'#2563eb'},'view':{'snapToObjects':True,'snapToGrid':False,'showGrid':False,'showGuides':True,'showElementLabels':False,'gridSize':2.5,'guideX':[50.0],'guideY':[50.0]},'slides':slides}
    outdeck.parent.mkdir(parents=True,exist_ok=True);outdeck.write_text(json.dumps(deck,ensure_ascii=False,indent=2)+'\n',encoding='utf-8');print(outdeck)
if __name__=='__main__':main()
