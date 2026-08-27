#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=1.0.2", "Pillow>=10"]
# ///
from __future__ import annotations
import io,json,math,sys
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE,MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN,MSO_VERTICAL_ANCHOR,MSO_AUTO_SIZE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.xmlchemy import OxmlElement

SW,SH=13.333333,7.5

def inchx(v):return Inches(float(v)*SW/100)
def inchy(v):return Inches(float(v)*SH/100)
def rgb(v,default='000000'):
    try:return RGBColor.from_string(v.lstrip('#')) if isinstance(v,str) and len(v.lstrip('#'))==6 else RGBColor.from_string(default)
    except:return RGBColor.from_string(default)
def set_arrowheads(shape, mode):
    if mode not in ('end','both'): return
    ln=shape._element.spPr.get_or_add_ln()
    for tag in (('a:headEnd', mode=='both'),('a:tailEnd', True)):
        if not tag[1]: continue
        node=OxmlElement(tag[0]);node.set('type','triangle');node.set('w','med');node.set('len','med');ln.append(node)


def set_text(box,text,e):
    tf=box.text_frame;tf.clear();tf.word_wrap=True;tf.vertical_anchor={'top':MSO_VERTICAL_ANCHOR.TOP,'middle':MSO_VERTICAL_ANCHOR.MIDDLE,'bottom':MSO_VERTICAL_ANCHOR.BOTTOM}.get(e.get('verticalAlign'),MSO_VERTICAL_ANCHOR.TOP);tf.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if e.get('autoFit') else MSO_AUTO_SIZE.NONE
    lines=str(text).split('\n') if str(text) else ['']
    for i,line in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.text=line;p.alignment={'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}.get(e.get('align'),PP_ALIGN.LEFT);p.line_spacing=float(e.get('lineSpacing',1.2));p.space_after=Pt(float(e.get('paragraphSpacing',0)))
        ppr=p._p.get_or_add_pPr();indent=float(e.get('indent',0));
        if indent>0:ppr.set('marL',str(int(Pt(indent))))
        ls=e.get('listStyle','none')
        if ls=='bullet':node=OxmlElement('a:buChar');node.set('char','•');ppr.append(node)
        elif ls=='number':node=OxmlElement('a:buAutoNum');node.set('type','arabicPeriod');ppr.append(node)
        for run in p.runs:run.font.name=e.get('fontFamily','Arial');run.font.size=Pt(float(e.get('fontSize',24)));run.font.bold=e.get('fontWeight',400)>=700;run.font.italic=bool(e.get('italic'));run.font.underline=bool(e.get('underline'));run.font.color.rgb=rgb(e.get('color'),'202124')

def set_flip(shape,e):
    if not e.get('flipH') and not e.get('flipV'):return
    try:
        xfrm=shape._element.spPr.xfrm
        if e.get('flipH'):xfrm.set('flipH','1')
        if e.get('flipV'):xfrm.set('flipV','1')
    except Exception:pass

def fill_line(shape,e):
    fill=e.get('fill','transparent');
    if fill=='transparent': shape.fill.background()
    else: shape.fill.solid(); shape.fill.fore_color.rgb=rgb(fill,'FFFFFF')
    stroke=e.get('stroke','transparent');
    if stroke=='transparent' or float(e.get('strokeWidth',1))<=0: shape.line.fill.background()
    else: shape.line.color.rgb=rgb(stroke,'4B4D50'); shape.line.width=Pt(float(e.get('strokeWidth',1)))
def cropped_image(deck_path,e):
    src=(deck_path.parent/e['path']).resolve(); im=Image.open(src); cr=e.get('crop',{}); l,t,r,b=(float(cr.get(k,0)) for k in ('left','top','right','bottom')); iw,ih=im.size; box=(int(iw*l/100),int(ih*t/100),max(1,int(iw*(100-r)/100)),max(1,int(ih*(100-b)/100))); im=im.crop(box); buf=io.BytesIO(); im.save(buf,'PNG'); buf.seek(0); return buf

def endpoint(ep,other,byid):
    e=byid.get(ep.get('elementId'))
    if not e:return(float(ep.get('x',10)),float(ep.get('y',10)))
    cx=e.get('x',0)+e.get('width',0)/2;cy=e.get('y',0)+e.get('height',0)/2;a=ep.get('anchor','auto')
    if a=='auto':a=('right' if other[0]>=cx else 'left') if abs(other[0]-cx)>abs(other[1]-cy) else ('bottom' if other[1]>=cy else 'top')
    return {'top':(cx,e['y']),'right':(e['x']+e['width'],cy),'bottom':(cx,e['y']+e['height']),'left':(e['x'],cy),'center':(cx,cy)}.get(a,(cx,cy))
def connector_points(e,byid):
    def center(ep,d):
        x=byid.get(ep.get('elementId'));return(x['x']+x['width']/2,x['y']+x['height']/2) if x else d
    a=center(e['from'],(10,10));b=center(e['to'],(80,80));a=endpoint(e['from'],b,byid);b=endpoint(e['to'],a,byid);return a,b

def add_element(slide,e,deck_path,byid):
    if e['type']=='connector':
        a,b=connector_points(e,byid); sh=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,inchx(a[0]),inchy(a[1]),inchx(b[0]),inchy(b[1])); sh.line.color.rgb=rgb(e.get('stroke'),'4B4D50');sh.line.width=Pt(float(e.get('strokeWidth',2)));
        if e.get('dash'): sh.line.dash_style=MSO_LINE_DASH_STYLE.DASH
        set_arrowheads(sh,e.get('arrow','end'));return
    x,y,w,h=map(float,(e.get('x',0),e.get('y',0),e.get('width',10),e.get('height',10))); left,top,width,height=inchx(x),inchy(y),inchx(w),inchy(h);typ=e['type']
    if typ=='image':
        sh=slide.shapes.add_picture(cropped_image(deck_path,e),left,top,width,height);sh.rotation=float(e.get('rotation',0));set_flip(sh,e);return
    if typ=='text':
        sh=slide.shapes.add_textbox(left,top,width,height);set_text(sh,e.get('text',''),e);fill_line(sh,e);sh.rotation=float(e.get('rotation',0));
        if e.get('role') in ('title','body'): sh.name=f"MorrowPresenter:{e['role']}"
        return
    if typ=='shape':
        st={'rect':MSO_SHAPE.RECTANGLE,'rounded-rect':MSO_SHAPE.ROUNDED_RECTANGLE,'ellipse':MSO_SHAPE.OVAL,'triangle':MSO_SHAPE.ISOSCELES_TRIANGLE,'diamond':MSO_SHAPE.DIAMOND,'pentagon':MSO_SHAPE.PENTAGON,'hexagon':MSO_SHAPE.HEXAGON,'star':MSO_SHAPE.STAR_5_POINT,'chevron':MSO_SHAPE.CHEVRON}.get(e.get('shape'))
        if e.get('shape') in ('line','arrow'):
            sh=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,left,top+height//2,left+width,top+height//2);sh.line.color.rgb=rgb(e.get('stroke'),'4B4D50');sh.line.width=Pt(float(e.get('strokeWidth',1.5)));set_arrowheads(sh,'end' if e.get('shape')=='arrow' else 'none');return
        sh=slide.shapes.add_shape(st or MSO_SHAPE.RECTANGLE,left,top,width,height);fill_line(sh,e);sh.rotation=float(e.get('rotation',0));set_flip(sh,e);set_text(sh,e.get('text',''),e);return
    if typ=='table':
        sh=slide.shapes.add_table(e['rows'],e['cols'],left,top,width,height);table=sh.table
        for r in range(e['rows']):
            for c in range(e['cols']):
                cell=table.cell(r,c);cell.text=e['cells'][r][c];cell.fill.solid();cell.fill.fore_color.rgb=rgb(e.get('headerFill') if r==0 else e.get('fill'),'FFFFFF');
                for p in cell.text_frame.paragraphs:
                    p.alignment={'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}.get(e.get('align'),PP_ALIGN.LEFT)
                    for run in p.runs:run.font.name=e.get('fontFamily','Arial');run.font.size=Pt(float(e.get('fontSize',18)));run.font.color.rgb=rgb(e.get('color'),'202124')

def add_footer(slide,deck,index):
    f=deck.get('footer',{});
    if f.get('showText'):
        sh=slide.shapes.add_textbox(inchx(4),inchy(94),inchx(72),inchy(4));sh.name='MorrowPresenter:footer';set_text(sh,f.get('text',''),{'fontFamily':f.get('fontFamily','Inter'),'fontSize':f.get('fontSize',11),'fontWeight':400,'italic':False,'underline':False,'color':f.get('color','#666666'),'align':'left','verticalAlign':'bottom','listStyle':'none','lineSpacing':1,'paragraphSpacing':0,'indent':0,'autoFit':False});sh.fill.background();sh.line.fill.background()
    if f.get('showSlideNumber'):
        sh=slide.shapes.add_textbox(inchx(88),inchy(94),inchx(8),inchy(4));sh.name='MorrowPresenter:slide-number';set_text(sh,str(index+1),{'fontFamily':f.get('fontFamily','Inter'),'fontSize':f.get('fontSize',11),'fontWeight':400,'italic':False,'underline':False,'color':f.get('color','#666666'),'align':'right','verticalAlign':'bottom','listStyle':'none','lineSpacing':1,'paragraphSpacing':0,'indent':0,'autoFit':False});sh.fill.background();sh.line.fill.background()

def main():
    if len(sys.argv)!=3:raise SystemExit('usage: export-pptx.py deck.morrowdeck output.pptx')
    deck_path=Path(sys.argv[1]).resolve();out=Path(sys.argv[2]).resolve();deck=json.loads(deck_path.read_text());page=deck.get('page',{});globals()['SW']=float(page.get('width',13.333333));globals()['SH']=float(page.get('height',7.5));prs=Presentation();prs.slide_width=Inches(SW);prs.slide_height=Inches(SH);blank=prs.slide_layouts[6]
    # remove default first slide if any only as created by template no slides normally
    for index,data in enumerate(deck['slides']):
        slide=prs.slides.add_slide(blank); bg=slide.background.fill;bg.solid();bg.fore_color.rgb=rgb(data.get('background',deck.get('theme',{}).get('background')),'FFFFFF')
        byid={e['id']:e for e in data.get('elements',[])}
        for e in data.get('elements',[]):add_element(slide,e,deck_path,byid)
        add_footer(slide,deck,index)
        if data.get('notes'):
            try: slide.notes_slide.notes_text_frame.text=data['notes']
            except Exception: pass
    out.parent.mkdir(parents=True,exist_ok=True);prs.save(out);print(out)
if __name__=='__main__':main()
